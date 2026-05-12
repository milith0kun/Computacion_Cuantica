"""
Exporta presentacion.html a un archivo .pptx (PowerPoint) editable.

Estrategia:
  1. Inicia un servidor HTTP local para que Playwright cargue el HTML con sus
     imágenes (fuentes Google se piden por internet).
  2. Abre el HTML en Chromium headless a 1920x1080 (formato PPT widescreen
     estándar 16:9).
  3. Para cada <section class="slide">, hace scroll a esa sección y captura
     un PNG a página completa.
  4. Compone un .pptx 16:9 donde cada slide del PPT corresponde a la captura
     de un slide del HTML.

Requisitos:
  - python-pptx   (ya instalado globalmente)
  - playwright    (se instala automáticamente si falta)

Uso:
    python exportar_pptx.py
"""

from __future__ import annotations

import http.server
import socketserver
import subprocess
import sys
import threading
import time
from pathlib import Path

# --------------------------------------------------------------------
# Configuración
# --------------------------------------------------------------------
HERE = Path(__file__).resolve().parent
HTML_FILE = HERE / "presentacion.html"
OUT_PPTX = HERE / "presentacion.pptx"
SHOTS_DIR = HERE / "_pptx_shots"
PORT = 8765
WIDTH = 1920
HEIGHT = 1080  # 16:9 widescreen estándar de PowerPoint


# --------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------
def ensure_pkg(import_name: str, pip_name: str | None = None) -> None:
    """Importa un paquete; si falla lo instala con pip."""
    pip_name = pip_name or import_name
    try:
        __import__(import_name)
    except ImportError:
        print(f"  -Instalando {pip_name} ...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "--quiet", pip_name])


def ensure_playwright_browser() -> None:
    """Descarga Chromium si Playwright aún no lo tiene."""
    try:
        # Intento rápido: si ya hay browser, no descargamos nada
        from playwright.sync_api import sync_playwright

        with sync_playwright() as p:
            browser = p.chromium.launch()
            browser.close()
        return
    except Exception:
        pass

    print("  -Descargando Chromium para Playwright (primera ejecución) ...")
    subprocess.check_call([sys.executable, "-m", "playwright", "install", "chromium"])


# --------------------------------------------------------------------
# Servidor local en hilo aparte
# --------------------------------------------------------------------
class QuietHandler(http.server.SimpleHTTPRequestHandler):
    """SimpleHTTPRequestHandler silencioso (sin spamear la consola)."""

    def log_message(self, *args, **kwargs):  # noqa: D401, ANN001, ANN002
        return


def start_server(directory: Path, port: int) -> socketserver.TCPServer:
    handler_factory = lambda *args, **kwargs: QuietHandler(  # noqa: E731
        *args, directory=str(directory), **kwargs
    )
    httpd = socketserver.TCPServer(("127.0.0.1", port), handler_factory)
    thread = threading.Thread(target=httpd.serve_forever, daemon=True)
    thread.start()
    return httpd


# --------------------------------------------------------------------
# Captura de slides
# --------------------------------------------------------------------
def capture_slides(url: str, out_dir: Path) -> list[Path]:
    """Abre el HTML, recorre cada `.slide` y guarda un PNG por slide."""
    from playwright.sync_api import sync_playwright

    out_dir.mkdir(exist_ok=True)
    # Limpiar capturas anteriores
    for old in out_dir.glob("slide_*.png"):
        old.unlink()

    shots: list[Path] = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        ctx = browser.new_context(
            viewport={"width": WIDTH, "height": HEIGHT},
            device_scale_factor=1,
            reduced_motion="reduce",
        )
        page = ctx.new_page()

        print("  -Cargando presentación ...")
        page.goto(url, wait_until="networkidle")

        # Esperar a que las fuentes web carguen del todo
        try:
            page.wait_for_function("document.fonts.ready.then(() => true)", timeout=8000)
        except Exception:
            pass
        time.sleep(0.6)  # respiro para que dibuje fondos/transiciones

        # Forzar a que todos los .reveal queden visibles (sin animación)
        page.evaluate(
            """
            document.querySelectorAll('.slide').forEach(s => s.classList.add('visible'));
            document.querySelectorAll('.reveal').forEach(el => {
                el.style.opacity = '1';
                el.style.transform = 'none';
                el.style.filter = 'none';
            });
            """
        )

        # Ocultar UI flotante (progress bar y nav dots) en las capturas
        page.add_style_tag(
            content="""
            .progress-bar, .nav-dots { display: none !important; }
            html { scroll-behavior: auto !important; }
            """
        )

        slides = page.query_selector_all(".slide")
        total = len(slides)
        print(f"  -{total} slides detectadas")

        for idx, slide in enumerate(slides, start=1):
            slide.scroll_into_view_if_needed()
            # Truco: scroll preciso al borde superior del slide
            page.evaluate(
                "(el) => el.scrollIntoView({behavior:'auto', block:'start'})", slide
            )
            page.wait_for_timeout(220)  # asentar render

            out = out_dir / f"slide_{idx:02d}.png"
            page.screenshot(path=str(out), clip={
                "x": 0,
                "y": 0,
                "width": WIDTH,
                "height": HEIGHT,
            }, type="png")
            shots.append(out)
            print(f"    [{idx:02d}/{total}] capturado")

        browser.close()

    return shots


# --------------------------------------------------------------------
# Composición del PPTX
# --------------------------------------------------------------------
def build_pptx(shots: list[Path], out_path: Path) -> None:
    """Construye un PPTX 16:9 colocando cada PNG como slide a página completa."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # PowerPoint widescreen 16:9: 13.333" x 7.5"
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # 'Blank'

    for shot in shots:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(shot),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(out_path))


# --------------------------------------------------------------------
# Main
# --------------------------------------------------------------------
def main() -> int:
    if not HTML_FILE.exists():
        print(f"[X] No se encontró {HTML_FILE}")
        return 1

    print("[*] Preparando dependencias")
    ensure_pkg("pptx", "python-pptx")
    ensure_pkg("playwright")
    ensure_playwright_browser()

    print(f"[*] Iniciando servidor local en :{PORT}")
    httpd = start_server(HERE, PORT)
    url = f"http://127.0.0.1:{PORT}/{HTML_FILE.name}"

    try:
        print("[*] Capturando slides a 1920x1080")
        shots = capture_slides(url, SHOTS_DIR)
    finally:
        httpd.shutdown()

    if not shots:
        print("[X] No se capturó ningún slide.")
        return 2

    print(f"[*] Ensamblando {OUT_PPTX.name}")
    build_pptx(shots, OUT_PPTX)

    print()
    print(f"[OK] Listo: {OUT_PPTX}")
    print(f"  - {len(shots)} slides · 13.333\" x 7.5\" (16:9)")
    print(f"  - Capturas intermedias: {SHOTS_DIR}")
    print()
    print("Abre el archivo con: PowerPoint, Keynote o LibreOffice Impress.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
